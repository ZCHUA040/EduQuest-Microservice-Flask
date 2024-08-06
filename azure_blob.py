from azure.storage.blob import BlobServiceClient
from docx import Document
from io import BytesIO
from pptx import Presentation
from pypdf import PdfReader


class AzureBlob:
    def __init__(self, connection_string, container_name):
        self.blob_service_client = BlobServiceClient.from_connection_string(connection_string)
        self.container_name = container_name

    def retrieve_document(self, document_id):
        """
        Retrieve a document from Azure Blob Storage
        :param document_id: The document ID
        :return: The document content in bytes
        """
        # Get the BlobClient
        blob_client = self.blob_service_client.get_blob_client(
            container=self.container_name,
            blob=document_id
        )

        # Download the blob content
        download_stream = blob_client.download_blob()
        document_content = download_stream.readall()

        if self.get_document_extension(document_id) == 'docx':
            return self.extract_text_from_docx(document_content)
        elif self.get_document_extension(document_id) == 'pdf':
            return self.extract_text_from_pdf(document_content)
        elif self.get_document_extension(document_id) == 'pptx':
            return self.extract_text_from_pptx(document_content)


    def get_document_extension(self, document_id):
        """
        Get the extension of a document
        :param document_id: The document ID
        :return: The extension of the document
        """
        return document_id.split('.')[-1]

    def extract_text_from_docx(self, document_bytes):
        """
        Extract text from a Word document
        :param document_bytes: The document in bytes
        :return: The extracted text in string format
        """
        # Read the docx file from the memory stream
        docx_stream = BytesIO(document_bytes)
        doc = Document(docx_stream)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)

    def extract_text_from_pdf(self, document_bytes):
        """
        Extract text from a PDF document
        :param document_bytes: The document in bytes
        :return: The extracted text in string format
        """
        pdf_stream = BytesIO(document_bytes)
        pdf = PdfReader(pdf_stream)
        full_text = []
        for page in pdf.pages:
            full_text.append(page.extract_text())
        return '\n'.join(full_text)

    def extract_text_from_pptx(self, document_bytes):
        """
        Extract text from a PowerPoint document
        :param document_bytes: The document in bytes
        :return: The extracted text in string format
        """
        pptx_stream = BytesIO(document_bytes)
        ppt = Presentation(pptx_stream)
        full_text = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        full_text.append(run.text)

        return '\n'.join(full_text)
